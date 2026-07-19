"""Read a PPTX into the IR.

Walks every shape including groups, keeps run and font data, speaker notes,
tables and images, and recovers SmartArt text. A run counts as code when its
font is monospace or it uses a command colour (see CMD_SCHEME_COLORS). Word-based
detection for the remaining cases happens later in classify.py.
"""

from __future__ import annotations

import hashlib
import re

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from .ir import Block, Deck, Image, Segment, Slide

_MONO_RE = re.compile(r"(mono|courier|consol|menlo|monaco|inconsolata|fira\s*code|"
                      r"source\s*code|dejavu\s*sans\s*mono|ibm\s*plex\s*mono|"
                      r"\bhack\b|jetbrains|cascadia|roboto\s*mono|ubuntu\s*mono|"
                      r"sf\s*mono|andale|lucida\s*console|pt\s*mono|space\s*mono|liberation\s*mono)", re.I)
_PAGENO_RE = re.compile(r"^\s*\d{1,3}\s*$")

# Some decks mark command lines with a theme colour (tx2, dk2 in LibreOffice).
# A run in one of these colours counts as a command. If a deck uses a fixed
# colour instead, add its hex value (upper-case, no '#') to CMD_RGB_COLORS.
CMD_SCHEME_COLORS = {"tx2", "dk2"}
CMD_RGB_COLORS: set[str] = set()

_A_SCHEMECLR = qn("a:schemeClr")
_A_SRGBCLR = qn("a:srgbClr")


def _is_mono(font_name: str | None) -> bool:
    return bool(font_name and _MONO_RE.search(font_name))


def _is_cmd_color(run) -> bool:
    """True when the run is explicitly painted in a configured command colour."""
    rpr = run._r.find(qn("a:rPr"))
    if rpr is None:
        return False
    for sc in rpr.iter(_A_SCHEMECLR):
        if sc.get("val") in CMD_SCHEME_COLORS:
            return True
    if CMD_RGB_COLORS:
        wanted = {c.upper() for c in CMD_RGB_COLORS}
        for rc in rpr.iter(_A_SRGBCLR):
            if (rc.get("val") or "").upper() in wanted:
                return True
    return False


def _para_segments(paragraph, font_code: bool = True, terms: bool = False) -> tuple[list[Segment], bool]:
    """Return (segments, all_literal) for one paragraph.

    A run painted in a command colour becomes a command. With terms on, a
    monospace run in the default colour becomes a term. Otherwise, when font_code
    is on, a monospace run becomes a command. all_literal is True when every
    non-blank run is command or term, which can promote the frame to a code figure.
    """
    segs: list[Segment] = []
    lit_flags: list[bool] = []
    for run in paragraph.runs:
        text = run.text
        if not text:
            continue
        mono = _is_mono(run.font.name)
        if _is_cmd_color(run):
            kind = "code"
        elif terms and mono:
            kind = "term"
        elif not terms and font_code and mono:
            kind = "code"
        else:
            kind = "prose"
        emph = bool(run.font.bold or run.font.italic)
        segs.append(Segment(text=text, kind=kind, emphasis=emph))
        if text.strip():
            lit_flags.append(kind in ("code", "term"))
    return segs, (bool(lit_flags) and all(lit_flags))


def _deck_uses_cmd_color(prs) -> bool:
    """True when any run in the deck is painted in a command colour."""
    def walk(shapes):
        for sh in shapes:
            try:
                if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                    if walk(sh.shapes):
                        return True
                    continue
            except Exception:
                pass
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip() and _is_cmd_color(run):
                            return True
        return False
    return any(walk(s.shapes) for s in prs.slides)


def _title_from_slide(slide) -> tuple[str, int | None]:
    """Title text and shape id from the title placeholder (p:ph type title/ctrTitle).

    Reads the placeholder by its type attribute, which is more reliable on messy
    decks than python-pptx's slide.shapes.title.
    """
    sps = slide._element.xpath(
        ".//p:sp[.//p:ph[@type='title'] or .//p:ph[@type='ctrTitle']]")
    if not sps:
        return "", None
    sp = sps[0]
    title = re.sub(r"\s+", " ", "".join(sp.xpath(".//a:t/text()"))).strip()
    ids = sp.xpath(".//p:cNvPr/@id")
    sid = int(ids[0]) if ids else None
    return title, sid

def _merge_prose(segs: list[Segment]) -> str:
    return "".join(s.text for s in segs)

def _extract_text_frame(tf, blocks: list[Block], font_code: bool = True,
                        terms: bool = False) -> None:
    """One source text frame -> exactly one block (analogous to one slide box).

    A frame whose paragraphs are all code and span >=2 lines becomes a code figure
    (e.g. a terminal dump). Everything else becomes a single `flow` block: a nested
    bullet list that keeps indent levels and shows inline commands in place.
    """
    lines: list[tuple[int, list[Segment], bool]] = []
    for paragraph in tf.paragraphs:
        segs, all_mono = _para_segments(paragraph, font_code, terms)
        if not _merge_prose(segs).strip():
            continue
        lines.append((paragraph.level or 0, segs, all_mono))
    if not lines:
        return
    if len(lines) >= 2 and all(m for _, _, m in lines):
        code = "\n".join(_merge_prose(segs) for _, segs, _ in lines)
        blocks.append(Block(kind="code", segments=[Segment(text=code, kind="code")]))
    else:
        blocks.append(Block(kind="flow", lines=[(lvl, segs) for lvl, segs, _ in lines]))

def _extract_table(shape, blocks: list[Block], font_code: bool = True,
                   terms: bool = False) -> None:
    rows: list[list[list[Segment]]] = []
    for r in shape.table.rows:
        row: list[list[Segment]] = []
        for cell in r.cells:
            cell_segs: list[Segment] = []
            for paragraph in cell.text_frame.paragraphs:
                segs, _ = _para_segments(paragraph, font_code, terms)
                if _merge_prose(segs).strip():
                    cell_segs.extend(segs)
                    cell_segs.append(Segment(text="\n", kind="prose"))
            if cell_segs and cell_segs[-1].text == "\n":
                cell_segs.pop()
            row.append(cell_segs or [Segment(text="", kind="prose")])
        rows.append(row)
    if rows:
        blocks.append(Block(kind="table", rows=rows))

_A_T = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"

def _drawingml_texts(shape) -> list[str]:
    """Best-effort: pull <a:t> text out of a shape's XML (recovers most SmartArt)."""
    try:
        return [e.text for e in shape._element.iter(_A_T) if e.text and e.text.strip()]
    except Exception:
        return []

def _extract_chart(shape, blocks: list[Block]) -> bool:
    """Best-effort chart text: title, category labels, series names."""
    try:
        chart = shape.chart
    except Exception:
        return False
    lines: list[tuple[int, list[Segment]]] = []
    try:
        if chart.has_title and chart.chart_title.text_frame.text.strip():
            lines.append((0, [Segment(text=chart.chart_title.text_frame.text.strip(), kind="prose")]))
    except Exception:
        pass
    try:
        cats = [str(c) for c in chart.plots[0].categories if str(c).strip()]
        for c in cats:
            lines.append((1, [Segment(text=c, kind="prose")]))
    except Exception:
        pass
    try:
        for s in chart.series:
            if getattr(s, "name", None):
                lines.append((1, [Segment(text=str(s.name), kind="prose")]))
    except Exception:
        pass
    if lines:
        blocks.append(Block(kind="flow", lines=lines))
        return True
    return False

def _walk(shapes, slide_idx: int, blocks: list[Block], images: list[Image],
          seen_hashes: dict[str, str], title_id, warnings: list[str],
          font_code: bool = True, terms: bool = False, order: str = "document") -> None:
    if order == "position":
        shapes = sorted(shapes, key=lambda sh: (sh.top if sh.top is not None else 10**9,
                                                sh.left if sh.left is not None else 0))
    for shape in shapes:
        if title_id is not None and shape.shape_id == title_id:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _walk(shape.shapes, slide_idx, blocks, images, seen_hashes, title_id,
                  warnings, font_code, terms, order)
            continue
        if shape.has_table:
            _extract_table(shape, blocks, font_code, terms)
            continue
        if getattr(shape, "has_chart", False):
            if not _extract_chart(shape, blocks):
                warnings.append(f"slide {slide_idx + 1}: chart text could not be read")
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
            try:
                blob = shape.image.blob
                h = hashlib.sha1(blob).hexdigest()
                if h not in seen_hashes:
                    fname = f"slide_{slide_idx + 1}_img_{len(seen_hashes) + 1}.{shape.image.ext}"
                    seen_hashes[h] = fname
                    images.append(Image(filename=fname, data=blob))
                continue
            except Exception:
                pass
        if shape.has_text_frame:
            txt = shape.text_frame.text
            if _PAGENO_RE.match(txt):
                continue
            _extract_text_frame(shape.text_frame, blocks, font_code, terms)
            continue
        if str(shape.shape_type) == "MEDIA (16)" or getattr(shape.shape_type, "name", "") == "MEDIA":
            warnings.append(f"slide {slide_idx + 1}: embedded media (video/audio) not imported")
            continue
        texts = _drawingml_texts(shape)
        if texts:
            blocks.append(Block(kind="flow",
                                lines=[(0, [Segment(text=t, kind="prose")]) for t in texts]))
            warnings.append(f"slide {slide_idx + 1}: diagram/object text imported flat "
                            f"(layout lost), please verify")

def _notes(slide) -> str:
    try:
        if slide.has_notes_slide:
            return (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:
        pass
    return ""

def extract(path: str, command_mode: str = "auto", terms: bool = False,
            order: str = "document") -> Deck:
    prs = Presentation(path)
    deck = Deck(title="", source_lang="en")
    # command detection: colour is always used; the monospace font is used unless
    # the deck marks commands by colour (auto), or the caller forces a mode.
    if command_mode == "color":
        font_code = False
    elif command_mode in ("font", "both"):
        font_code = True
    else:
        font_code = not _deck_uses_cmd_color(prs)
    seen_hashes: dict[str, str] = {}
    for i, slide in enumerate(prs.slides):
        # title from the title placeholder (by type), not python-pptx's guess
        title, title_id = _title_from_slide(slide)
        blocks: list[Block] = []
        images: list[Image] = []
        _walk(slide.shapes, i, blocks, images, seen_hashes, title_id, deck.warnings,
              font_code, terms, order)
        existing = " ".join(b.text() for b in blocks)
        for rel in slide.part.rels.values():
            try:
                if "diagramData" not in rel.reltype:
                    continue
                root = etree.fromstring(rel.target_part.blob)
                dtexts = [e.text.strip() for e in root.iter(_A_T)
                          if e.text and e.text.strip() and e.text.strip() not in existing]
                if dtexts:
                    blocks.append(Block(kind="flow",
                                        lines=[(0, [Segment(text=t, kind="prose")]) for t in dtexts]))
                    deck.warnings.append(f"slide {i + 1}: SmartArt text imported flat "
                                         f"(layout lost), please verify")
            except Exception:
                pass
        notes = _notes(slide)
        if not title and not blocks and not images:
            if not notes:
                continue
            title = notes.split("\n", 1)[0].strip()[:60]
            blocks.append(Block(kind="flow", lines=[
                (0, [Segment(text=ln.strip(), kind="prose")])
                for ln in notes.splitlines() if ln.strip()]))
            notes = ""
            deck.warnings.append(f"slide {i + 1}: only speaker notes present, imported as content")
        deck.slides.append(Slide(index=i, title=title, notes=notes,
                                 blocks=blocks, images=images))
    return deck
