"""Self-test / generalisation check over real decks.

Runs the whole pipeline (in memory, writes nothing) over one or more files or a
folder, and checks the things that must hold for *every* deck, not just the one it
was built on:

  * XML well-formed
  * fidelity: every source segment present, code byte-exact
  * no empty content/notes blocks (Scenari rejects empty mandatory fields)
  * every sc:refUri resolves to a produced file (no dangling references)
  * PPTX only: independent source coverage (all shapes + tables + notes vs the
    emitted XML), catches text the extractor drops (SmartArt, etc.)
  * element vocabulary stays within the known Dokiel set (built from reference
    material under archive/, e.g. the course skeleton; skipped if none is found)

Usage:
    python3 -m converter.selftest FILE_OR_DIR [MORE ...]
    python3 -m converter.selftest ~/decks            # recurse a folder
Exit code 0 if every file passes, 1 otherwise.
"""

from __future__ import annotations

import posixpath
import re
import sys
from pathlib import Path

from lxml import etree

from . import classify, cli, emit, verify

SUPPORTED = (".pptx", ".ppt", ".pdf", ".docx")
_REFURI = "{http://www.utc.fr/ics/scenari/v3/core}refUri"
_EMPTY_PATTERNS = (
    '<sc:para xml:space="preserve"></sc:para>', "<sc:para/>", "<dk:comment/>",
    "<dk:content></dk:content>", "<dk:content/>", "<dk:text></dk:text>",
)
_REPO = Path(__file__).resolve().parent.parent

def _empty_blocks(files: dict) -> int:
    return sum(v.count(pat) for k, v in files.items()
               if isinstance(v, str) and k.endswith((".scen", ".unit"))
               for pat in _EMPTY_PATTERNS)

def _dangling_refs(files: dict) -> list[str]:
    keys = set(files.keys())
    dirs = {posixpath.dirname(k) for k in keys}
    bad = []
    for rel, content in files.items():
        if not rel.endswith((".scen", ".unit", ".pub")):
            continue
        data = content.encode("utf-8") if isinstance(content, str) else content
        try:
            root = etree.fromstring(data)
        except Exception:
            continue
        for e in root.iter():
            ref = e.get(_REFURI)
            if not ref or ref.startswith("../&"):
                continue
            target = posixpath.normpath(posixpath.join(posixpath.dirname(rel), ref))
            if target in keys or target in dirs:
                continue
            bad.append(f"{rel} -> {ref}")
    return bad

def _pptx_source_text(path: str) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    parts: list[str] = []

    def walk(shapes):
        for sh in shapes:
            try:
                if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(sh.shapes); continue
            except Exception:
                pass
            try:
                if sh.has_table:
                    for r in sh.table.rows:
                        for c in r.cells:
                            parts.append(c.text)
                    continue
            except Exception:
                pass
            try:
                if sh.has_text_frame:
                    parts.append(sh.text_frame.text)
            except Exception:
                pass

    prs = Presentation(path)
    for sl in prs.slides:
        walk(sl.shapes)
        try:
            if sl.has_notes_slide:
                parts.append(sl.notes_slide.notes_text_frame.text)
        except Exception:
            pass
    return " ".join(parts)

def _coverage(path: str, emitted_ns: str) -> tuple[float, list[str]]:
    """PPTX only: fraction of source words (>=4 chars) present in emitted XML."""
    src = {w for w in re.findall(r"\S{4,}", _pptx_source_text(path))
           if not re.fullmatch(r"\d{1,3}", w)}
    if not src:
        return 100.0, []
    missing = [w for w in src if re.sub(r"\s+", "", w) not in emitted_ns]
    return 100.0 * (len(src) - len(missing)) / len(src), missing

def _reference_vocab() -> set | None:
    dirs = [d for d in (_REPO / "archive", _REPO / "Course_Skeleton") if d.exists()]
    if not dirs:
        return None
    voc: set = set()
    for d in dirs:
        for ext in ("*.scen", "*.unit", "*.pub"):
            for p in d.rglob(ext):
                try:
                    for e in etree.parse(str(p)).getroot().iter():
                        if isinstance(e.tag, str) and "}" in e.tag:
                            voc.add(e.tag.split("}", 1)[1])
                except Exception:
                    pass
    return voc or None

def _foreign(files: dict, vocab: set) -> set:
    used: set = set()
    for rel, content in files.items():
        if not rel.endswith((".scen", ".unit", ".pub")):
            continue
        data = content.encode("utf-8") if isinstance(content, str) else content
        try:
            for e in etree.fromstring(data).iter():
                if isinstance(e.tag, str) and "}" in e.tag:
                    used.add(e.tag.split("}", 1)[1])
        except Exception:
            pass
    return used - vocab

def check_file(path: str, vocab: set | None) -> dict:
    deck = cli.load_deck(path, log=lambda *a: None)
    classify.classify(deck)
    files = emit.build_workspace(deck, "Selftest", "Day 1", "1-1", Path(path).stem)

    wf = verify.check_wellformed(files)
    fid = verify.check_fidelity(deck, files)
    empty = _empty_blocks(files)
    dangling = _dangling_refs(files)
    foreign = _foreign(files, vocab) if vocab else set()

    cov, cov_missing = (100.0, [])
    if Path(path).suffix.lower() in (".pptx", ".ppt"):
        cov, cov_missing = _coverage(path, re.sub(r"\s+", "", verify._emitted_text(files)))

    fails = []
    if wf:
        fails.append(f"{len(wf)} malformed XML ({wf[0]})")
    if fid["missing_tokens"]:
        fails.append(f"{len(fid['missing_tokens'])} segments missing")
    if not fid["code_ok"]:
        fails.append(f"{len(fid['code_missing'])} code blocks changed")
    if empty:
        fails.append(f"{empty} empty blocks")
    if dangling:
        fails.append(f"{len(dangling)} dangling refs ({dangling[0]})")
    if foreign:
        fails.append(f"foreign elements {sorted(foreign)}")
    if cov < 99.5:
        fails.append(f"coverage {cov:.1f}% (e.g. {cov_missing[:4]})")

    return {"slides": len(deck.slides), "cov": cov, "warnings": deck.warnings,
            "fails": fails}

def collect(paths: list[str]) -> list[Path]:
    out: list[Path] = []
    for p in paths:
        pp = Path(p)
        if pp.is_dir():
            out += [q for q in sorted(pp.rglob("*"))
                    if q.suffix.lower() in SUPPORTED and not q.name.startswith("~$")]
        elif pp.suffix.lower() in SUPPORTED:
            out.append(pp)
    return out

def main(argv=None) -> int:
    argv = argv if argv is not None else sys.argv[1:]
    if not argv:
        print("usage: python3 -m converter.selftest FILE_OR_DIR [...]")
        return 2
    files = collect(argv)
    if not files:
        print("no .pptx/.pdf/.docx found")
        return 2
    vocab = _reference_vocab()
    print(f"self-test over {len(files)} file(s)"
          + ("" if vocab else "  [vocabulary check skipped: no reference found]"))
    n_fail = n_warn = 0
    for f in files:
        try:
            r = check_file(str(f), vocab)
        except Exception as e:
            n_fail += 1
            print(f"  FAIL {f.name[:44]:44} EXCEPTION {type(e).__name__}: {e}")
            continue
        ok = not r["fails"]
        n_fail += 0 if ok else 1
        n_warn += len(r["warnings"])
        tag = "ok  " if ok else "FAIL"
        print(f"  {tag} {f.name[:44]:44} slides={r['slides']:3} cov={r['cov']:5.1f}% "
              f"warn={len(r['warnings'])}")
        for msg in r["fails"]:
            print(f"       - {msg}")
        for w in r["warnings"][:4]:
            print(f"       ~ {w}")
    print(f"\n{'PASS' if not n_fail else 'FAIL'}: {len(files) - n_fail}/{len(files)} "
          f"files clean, {n_warn} warning(s) total")
    return 1 if n_fail else 0

if __name__ == "__main__":
    sys.exit(main())
